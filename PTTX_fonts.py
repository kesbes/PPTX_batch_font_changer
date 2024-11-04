import os
from pptx import Presentation

# Update this to your actual directory path CHANGE YOUR CODE HERE
directory_path = r"C:\Users\kesbes\Desktop\PPTX"

# Function to change fonts in a presentation
def change_font_to_arial(ppt_path):
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'#CHANGE YOUR CODE HERE
    # Save the updated presentation
    presentation.save(ppt_path)

# Check if the directory exists
if os.path.exists(directory_path):
    for filename in os.listdir(directory_path):
        if filename.endswith(".pptx"):
            change_font_to_arial(os.path.join(directory_path, filename))
            print(f"Font updated for {filename}")
else:
    print(f"The specified directory does not exist: {directory_path}")
